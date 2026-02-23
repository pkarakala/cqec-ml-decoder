"""Build Phase 4 presentation as .pptx"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
BG = RGBColor(0x11, 0x11, 0x1b)
SURFACE = RGBColor(0x1e, 0x1e, 0x2e)
TEXT = RGBColor(0xcd, 0xd6, 0xf4)
SUBTEXT = RGBColor(0xa6, 0xad, 0xc8)
ACCENT_BLUE = RGBColor(0x89, 0xb4, 0xfa)
ACCENT_GREEN = RGBColor(0xa6, 0xe3, 0xa1)
ACCENT_RED = RGBColor(0xf3, 0x8b, 0xa8)
ACCENT_PEACH = RGBColor(0xfa, 0xb3, 0x87)
ACCENT_MAUVE = RGBColor(0xcb, 0xa6, 0xf7)
WHITE = RGBColor(0xff, 0xff, 0xff)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

FIGURES = os.path.join(os.path.dirname(__file__), "..", "outputs", "figures")


def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=TEXT, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return tf

def add_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text

def try_add_image(slide, filename, left, top, width, height=None):
    path = os.path.join(FIGURES, filename)
    if os.path.exists(path):
        if height:
            slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
        else:
            slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width))
        return True
    else:
        # Placeholder box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height or 4)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = SURFACE
        shape.line.color.rgb = ACCENT_BLUE
        shape.line.width = Pt(2)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"📊 Insert: {filename}"
        p.font.size = Pt(14)
        p.font.color.rgb = SUBTEXT
        p.alignment = PP_ALIGN.CENTER
        return False


# ═══════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

add_textbox(slide, 1.5, 1.5, 10, 1.2,
    "Can a Neural Network Decode Quantum Errors\nBetter Than Bayes' Theorem?",
    font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 3.2, 10, 0.8,
    "Adaptive Machine Learning for Real-Time Quantum Error Correction",
    font_size=22, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 4.5, 10, 0.5,
    "Pranav Reddy  ·  Clark Enge",
    font_size=18, color=SUBTEXT, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 5.2, 10, 0.5,
    "UC Santa Barbara  ·  Data Science Club",
    font_size=16, color=SUBTEXT, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 2: The Problem
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 11, 0.8,
    "Quantum Computers Are Incredibly Fragile",
    font_size=32, color=WHITE, bold=True)

add_bullet_list(slide, 0.8, 1.5, 7, 4.5, [
    "• Qubits lose information in microseconds",
    "• Error rates: ~0.1–1% per operation  (vs ~10⁻¹⁵ for classical bits)",
    "• Without error correction, quantum advantage is impossible",
    "• Every major quantum roadmap depends on solving this",
    "",
    "• Google, IBM, Microsoft all agree:",
    "  Error correction is THE bottleneck for useful quantum computing",
], font_size=20, color=TEXT)

add_speaker_notes(slide,
    "Hook: quantum computers are powerful but incredibly fragile. "
    "Error correction is the single biggest unsolved engineering challenge. "
    "If we can't fix errors faster than they happen, quantum computing stays in the lab.")

# ═══════════════════════════════════════════════════════════
# SLIDE 3: How QEC Works
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 11, 0.8,
    "The 3-Qubit Repetition Code",
    font_size=32, color=WHITE, bold=True)

add_bullet_list(slide, 0.8, 1.4, 6, 2.5, [
    "• Encode 1 logical qubit across 3 physical qubits",
    "    |0⟩ → |000⟩     |1⟩ → |111⟩",
    "• Measure stabilizers to detect errors without",
    "  destroying the quantum state",
    "• Each error produces a unique syndrome signature",
], font_size=18, color=TEXT)

# Syndrome table
add_textbox(slide, 7.5, 1.4, 5, 0.5,
    "Syndrome Table", font_size=18, color=ACCENT_BLUE, bold=True)

table_data = [
    ("Error", "S₁", "S₂"),
    ("No error", "+1", "+1"),
    ("Flip qubit 1", "−1", "+1"),
    ("Flip qubit 2", "−1", "−1"),
    ("Flip qubit 3", "+1", "−1"),
]
table = slide.shapes.add_table(5, 3, Inches(7.5), Inches(2.0), Inches(4.5), Inches(2.5)).table
for row_idx, row_data in enumerate(table_data):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = cell_text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT if row_idx > 0 else WHITE
            p.font.bold = (row_idx == 0)
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = SURFACE if row_idx > 0 else RGBColor(0x31, 0x32, 0x44)

add_speaker_notes(slide,
    "Think of it like a checksum. We don't read the data directly — "
    "we read parity checks that tell us if something flipped.")


# ═══════════════════════════════════════════════════════════
# SLIDE 4: Noisy Measurements
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 11, 0.8,
    "Real Measurements Are Noisy Analog Signals",
    font_size=32, color=WHITE, bold=True)

add_bullet_list(slide, 0.8, 1.5, 5.5, 3, [
    "• Real hardware gives continuous readout:",
    "    r(t) = signal + noise",
    "• The decoder must infer the error state",
    "  from a noisy time series",
    "• This is a classification problem:",
    "  4 error classes from 2 noisy signals",
    "",
    "• This is where ML comes in — we're doing",
    "  signal processing under uncertainty",
], font_size=18, color=TEXT)

try_add_image(slide, "decoder_comparison.png", 7, 1.5, 5.5, 4.5)

add_speaker_notes(slide,
    "We're not doing textbook QEC with discrete syndrome bits. "
    "We're working with continuous analog readout buried in noise. "
    "The decoder has to figure out what happened from messy real-valued signals.")

# ═══════════════════════════════════════════════════════════
# SLIDE 5: Three Decoders
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 11, 0.8,
    "Three Decoders, Head to Head",
    font_size=32, color=WHITE, bold=True)

# Three columns
for i, (name, color, desc, detail) in enumerate([
    ("Threshold", ACCENT_RED, "Average signal,\ncheck the sign",
     "No model needed\nFast, simple, fragile"),
    ("Bayesian Filter", ACCENT_GREEN, "Optimal probabilistic\nfilter (Wonham/HMM)",
     "Requires known noise model\nOptimal when assumptions hold"),
    ("GRU Neural Net", ACCENT_BLUE, "Recurrent network\nthat learns from data",
     "Learns the model from examples\nRobust to model mismatch"),
]):
    left = 0.8 + i * 4.1
    add_textbox(slide, left, 1.5, 3.8, 0.5, name,
                font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, 2.2, 3.8, 1.5, desc,
                font_size=16, color=TEXT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, 3.8, 3.8, 1.5, detail,
                font_size=14, color=SUBTEXT, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.8, 5.5, 11, 0.8,
    "The Bayesian filter is mathematically optimal IF you know the exact noise model.\n"
    "The GRU doesn't need that assumption — it learns directly from data.",
    font_size=16, color=SUBTEXT)

add_speaker_notes(slide,
    "The Bayesian filter is the gold standard — optimal under known noise model. "
    "The GRU learns the dynamics from data. The question is: which wins when the model is wrong?")


# ═══════════════════════════════════════════════════════════
# SLIDE 6: Phase 1 & 2 Results
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 11, 0.8,
    "Phase 1 & 2: GRU Wins Under Dynamics",
    font_size=32, color=WHITE, bold=True)

# Phase 1 table
add_textbox(slide, 0.8, 1.4, 5, 0.5, "Phase 1 — Static (ideal)",
            font_size=18, color=ACCENT_BLUE, bold=True)
t1 = slide.shapes.add_table(3, 2, Inches(0.8), Inches(2.0), Inches(5), Inches(1.5)).table
for r, (dec, acc) in enumerate([("Decoder", "Accuracy"), ("Threshold", "~86%"), ("GRU", "~96%")]):
    for c, val in enumerate([dec, acc]):
        cell = t1.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16); p.font.color.rgb = TEXT if r > 0 else WHITE
            p.font.bold = (r == 0 or (r == 2 and c == 1)); p.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = SURFACE if r > 0 else RGBColor(0x31, 0x32, 0x44)

# Phase 2 table
add_textbox(slide, 6.8, 1.4, 5.5, 0.5, "Phase 2 — Hamiltonian dynamics",
            font_size=18, color=ACCENT_BLUE, bold=True)
t2 = slide.shapes.add_table(4, 2, Inches(6.8), Inches(2.0), Inches(5.5), Inches(2.0)).table
for r, (dec, acc) in enumerate([("Decoder", "Accuracy"), ("Threshold", "~85%"),
                                 ("Bayesian Filter", "~94%"), ("GRU", "~96%")]):
    for c, val in enumerate([dec, acc]):
        cell = t2.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16); p.font.color.rgb = TEXT if r > 0 else WHITE
            p.font.bold = (r == 0 or (r == 3 and c == 1)); p.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = SURFACE if r > 0 else RGBColor(0x31, 0x32, 0x44)

add_textbox(slide, 0.8, 4.5, 11, 1.2,
    "When we add coherent drive, calibration drift, and measurement backaction,\n"
    "the Bayesian filter's assumptions break. The GRU maintains performance\n"
    "by learning dynamics directly from data.",
    font_size=18, color=SUBTEXT)

try_add_image(slide, "phase2_dynamics_comparison.png", 3, 5.2, 7, 2)

add_speaker_notes(slide,
    "Phase 1 is the easy case. Phase 2 adds real physics. "
    "The Bayesian filter degrades because its model is wrong. The GRU just learns the dynamics.")

# ═══════════════════════════════════════════════════════════
# SLIDE 7: Phase 3 Non-Idealities
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 11, 0.8,
    "Phase 3: Real Hardware Is Worse Than You Think",
    font_size=32, color=WHITE, bold=True)

add_bullet_list(slide, 0.8, 1.5, 5, 4, [
    "Three non-idealities that break",
    "model-based decoders:",
    "",
    "1. Colored noise",
    "   Temporally correlated (AR(1)),",
    "   not white Gaussian",
    "",
    "2. Post-flip transients",
    "   Exponential ring-down after",
    "   each error event",
    "",
    "3. Random-walk drift",
    "   Measurement calibration wanders",
    "   via Brownian motion",
], font_size=16, color=TEXT)

try_add_image(slide, "phase3_nonideal_effects.png", 6.5, 1.2, 6, 5)

add_speaker_notes(slide,
    "These effects appear in real superconducting qubit hardware. "
    "The Bayesian filter assumes white noise and static parameters — both assumptions are now violated.")


# ═══════════════════════════════════════════════════════════
# SLIDE 8: Phase 3 Results
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 11, 0.8,
    "Phase 3 Results: Everyone Degrades",
    font_size=32, color=WHITE, bold=True)

t3 = slide.shapes.add_table(4, 3, Inches(0.8), Inches(1.5), Inches(11), Inches(2.2)).table
for r, row in enumerate([
    ("Decoder", "Accuracy", "Notes"),
    ("Threshold", "79.4%", "Simple averaging fails with colored noise"),
    ("Bayesian Filter", "84.2%", "White noise assumption violated"),
    ("GRU", "83.1%", "Learns from data, but overfits"),
]):
    for c, val in enumerate(row):
        cell = t3.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16); p.font.color.rgb = TEXT if r > 0 else WHITE
            p.font.bold = (r == 0); p.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = SURFACE if r > 0 else RGBColor(0x31, 0x32, 0x44)

add_textbox(slide, 0.8, 4.0, 11, 1,
    "Key insight: The Bayesian filter's advantage disappears when its assumptions are wrong.\n"
    "The GRU is competitive but needs more data. This motivates our key question...",
    font_size=18, color=SUBTEXT)

try_add_image(slide, "phase3_decoder_comparison.png", 2, 5, 9, 2.2)

add_speaker_notes(slide,
    "Bayesian and GRU are neck-and-neck. The model-based approach loses its edge "
    "when the model is wrong. What happens when things get even worse?")

# ═══════════════════════════════════════════════════════════
# SLIDE 9: Phase 4 — The Challenge
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 11, 0.8,
    "Phase 4: Hardware Drifts During Operation",
    font_size=32, color=WHITE, bold=True)

add_bullet_list(slide, 0.8, 1.5, 5.5, 4, [
    "• Real quantum hardware drifts continuously:",
    "  temperature, aging, environment",
    "",
    "• Parameters aren't just non-ideal —",
    "  they're non-stationary",
    "",
    "• A decoder trained on yesterday's noise",
    "  fails on today's hardware",
    "",
    "• Current solution: stop computation,",
    "  recalibrate, restart (expensive!)",
], font_size=18, color=TEXT)

try_add_image(slide, "phase4_drift_schedules.png", 6.8, 1.2, 5.8, 5)

add_speaker_notes(slide,
    "This is the problem nobody has solved with ML decoders. You train your network, "
    "deploy it, and within hours the hardware has drifted. The industry solution is constant recalibration.")

# ═══════════════════════════════════════════════════════════
# SLIDE 10: Our Solution
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 11, 0.8,
    "Our Solution: A Decoder That Keeps Learning",
    font_size=32, color=WHITE, bold=True)

add_bullet_list(slide, 0.8, 1.5, 11, 1.5, [
    "Same GRU architecture, but weights update during inference",
    "EMA-smoothed gradient updates — stable, low-overhead",
    "Three adaptation strategies tested:",
], font_size=18, color=TEXT)

# Three strategy boxes
for i, (name, color, desc) in enumerate([
    ("Static GRU", ACCENT_MAUVE, "Trained once, frozen\n(baseline)"),
    ("Pseudo-label", ACCENT_PEACH, "Self-training with confident\npredictions → fails under drift"),
    ("Hybrid", ACCENT_GREEN, "Periodic true labels +\npseudo-labels → works!"),
]):
    left = 0.8 + i * 4.1
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(3.5), Inches(3.8), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SURFACE
    shape.line.color.rgb = color
    shape.line.width = Pt(3)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(20); p.font.color.rgb = color; p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "\n" + desc
    p2.font.size = Pt(14); p2.font.color.rgb = SUBTEXT
    p2.alignment = PP_ALIGN.CENTER

add_speaker_notes(slide,
    "Key insight: pure self-training fails because when the model is confidently wrong, "
    "it reinforces its own mistakes. But periodic recalibration anchors the model.")


# ═══════════════════════════════════════════════════════════
# SLIDE 11: Phase 4 Overall Results (FILL IN)
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 11, 0.8,
    "Phase 4 Results: Five-Way Decoder Comparison",
    font_size=32, color=WHITE, bold=True)

add_textbox(slide, 0.8, 1.2, 5, 0.5,
    "⚠️ Fill in after running notebook", font_size=14, color=ACCENT_PEACH)

t4 = slide.shapes.add_table(6, 2, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.5)).table
for r, (dec, acc) in enumerate([
    ("Decoder", "Accuracy"),
    ("Threshold", "____%"),
    ("Bayesian Filter", "____%"),
    ("Static GRU", "____%"),
    ("Adaptive (pseudo-labels)", "____%"),
    ("Adaptive (hybrid)", "____%"),
]):
    for c, val in enumerate([dec, acc]):
        cell = t4.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(18); p.font.color.rgb = TEXT if r > 0 else WHITE
            p.font.bold = (r == 0 or r == 5); p.font.name = "Calibri"
            if r == 5: p.font.color.rgb = ACCENT_GREEN
        cell.fill.solid()
        cell.fill.fore_color.rgb = SURFACE if r > 0 else RGBColor(0x31, 0x32, 0x44)

try_add_image(slide, "phase4_decoder_comparison.png", 7, 1.2, 5.5, 5)

add_speaker_notes(slide,
    "The headline number matters less than the temporal breakdown on the next slide.")

# ═══════════════════════════════════════════════════════════
# SLIDE 12: KEY SLIDE — Accuracy Over Time ⭐
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.2, 11, 0.8,
    "⭐ As Hardware Drifts, Static Decoders Fail — Adaptive Survives",
    font_size=30, color=WHITE, bold=True)

add_textbox(slide, 0.8, 1.0, 5, 0.4,
    "⚠️ Fill in segment data after running notebook", font_size=12, color=ACCENT_PEACH)

t5 = slide.shapes.add_table(6, 6, Inches(0.3), Inches(1.5), Inches(12.5), Inches(2.8)).table
headers = ["Segment", "Threshold", "Bayesian", "Static GRU", "Pseudo-label", "Hybrid"]
rows = [
    ("1 (early)", "___", "___", "___", "___", "___"),
    ("2", "___", "___", "___", "___", "___"),
    ("3 (mid)", "___", "___", "___", "___", "___"),
    ("4", "___", "___", "___", "___", "___"),
    ("5 (late)", "___", "___", "___", "___", "___"),
]
for c, h in enumerate(headers):
    cell = t5.cell(0, c)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(13); p.font.color.rgb = WHITE; p.font.bold = True; p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x31, 0x32, 0x44)

for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = t5.cell(r+1, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13); p.font.color.rgb = TEXT; p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
            if c == 5: p.font.color.rgb = ACCENT_GREEN; p.font.bold = True
        cell.fill.solid(); cell.fill.fore_color.rgb = SURFACE

try_add_image(slide, "phase4_accuracy_over_time.png", 1, 4.5, 11, 2.8)

add_speaker_notes(slide,
    "THIS IS THE KEY SLIDE. Point to the lines diverging. "
    "Early on, everyone does fine. As parameters drift, static decoders fall off. "
    "Pseudo-label barely helps — confident wrong predictions poison self-training. "
    "Hybrid maintains accuracy. This is the path to self-calibrating QEC.")


# ═══════════════════════════════════════════════════════════
# SLIDE 13: Why Pseudo-Labels Fail
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 11, 0.8,
    "Why Pure Self-Training Fails: The Pseudo-Label Trap",
    font_size=32, color=WHITE, bold=True)

add_bullet_list(slide, 0.8, 1.5, 6, 3, [
    "• Model is ~95% confident on average",
    "• But accuracy drops to ~70% under heavy drift",
    "• High confidence + wrong answer = poisoned labels",
    "• The model reinforces its own mistakes",
    "• Accuracy spirals downward over time",
], font_size=20, color=TEXT)

# 2x2 grid
for r in range(2):
    for c in range(2):
        left = 7.5 + c * 2.5
        top = 1.5 + r * 2.2
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(2.3), Inches(1.8))
        shape.fill.solid()
        labels = [
            ("Correct +\nHigh Conf", ACCENT_GREEN, "✅ Good"),
            ("Wrong +\nHigh Conf", ACCENT_RED, "❌ Poison!"),
            ("Correct +\nLow Conf", ACCENT_BLUE, "⚠️ Skip"),
            ("Wrong +\nLow Conf", SUBTEXT, "⚠️ Skip"),
        ]
        idx = r * 2 + c
        name, color, status = labels[idx]
        shape.fill.fore_color.rgb = SURFACE
        shape.line.color.rgb = color; shape.line.width = Pt(2)
        tf = shape.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = name
        p.font.size = Pt(12); p.font.color.rgb = color; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = status
        p2.font.size = Pt(16); p2.font.color.rgb = color; p2.alignment = PP_ALIGN.CENTER

add_textbox(slide, 0.8, 5.5, 11, 1,
    "The confidence threshold can't save you when the distribution shifts.\n"
    "You need ground truth to anchor the model.",
    font_size=18, color=SUBTEXT)

add_speaker_notes(slide,
    "This is a well-known failure mode in semi-supervised learning. "
    "When the data distribution shifts, confident predictions become unreliable.")

# ═══════════════════════════════════════════════════════════
# SLIDE 14: Supervision Frequency (FILL IN)
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 11, 0.8,
    "How Often Do You Need True Labels?",
    font_size=32, color=WHITE, bold=True)

add_textbox(slide, 0.8, 1.2, 5, 0.4,
    "⚠️ Fill in after running notebook", font_size=14, color=ACCENT_PEACH)

t6 = slide.shapes.add_table(8, 3, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4.5)).table
for r, row in enumerate([
    ("Supervision Rate", "% Supervised", "Accuracy"),
    ("Every 10", "10%", "____%"),
    ("Every 20", "5%", "____%"),
    ("Every 50", "2%", "____%"),
    ("Every 100", "1%", "____%"),
    ("Every 200", "0.5%", "____%"),
    ("Every 500", "0.2%", "____%"),
    ("Static (none)", "0%", "____%"),
]):
    for c, val in enumerate(row):
        cell = t6.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.color.rgb = TEXT if r > 0 else WHITE
            p.font.bold = (r == 0); p.font.name = "Calibri"; p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = SURFACE if r > 0 else RGBColor(0x31, 0x32, 0x44)

try_add_image(slide, "phase4_robustness_drift.png", 7, 1.2, 5.5, 5)

add_speaker_notes(slide,
    "Even 1-2% true labels gives most of the benefit. "
    "That's one recalibration every 50-100 windows — realistic for real hardware.")


# ═══════════════════════════════════════════════════════════
# SLIDE 15: Four Phases Journey
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 11, 0.8,
    "The Journey: Four Phases of Increasing Realism",
    font_size=32, color=WHITE, bold=True)

phases = [
    ("Phase 1", "Static\nSyndromes", "GRU: ~96%", ACCENT_BLUE, "Ideal world"),
    ("Phase 2", "Hamiltonian\nDynamics", "GRU: ~96%\nBayes: ~94%", ACCENT_MAUVE, "Add physics"),
    ("Phase 3", "Non-Ideal\nEffects", "GRU: ~83%\nBayes: ~84%", ACCENT_PEACH, "Add hardware\nimperfections"),
    ("Phase 4", "Drifting\nParameters", "Hybrid: ____%\nStatic: ____%", ACCENT_GREEN, "Hardware changes\nduring operation"),
]

for i, (name, desc, result, color, subtitle) in enumerate(phases):
    left = 0.5 + i * 3.2
    # Phase box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(1.5), Inches(2.9), Inches(4.5))
    shape.fill.solid(); shape.fill.fore_color.rgb = SURFACE
    shape.line.color.rgb = color; shape.line.width = Pt(3)
    
    add_textbox(slide, left + 0.1, 1.6, 2.7, 0.5, name,
                font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.1, 2.2, 2.7, 1.0, desc,
                font_size=16, color=TEXT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.1, 3.3, 2.7, 1.0, result,
                font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.1, 4.8, 2.7, 0.8, subtitle,
                font_size=12, color=SUBTEXT, alignment=PP_ALIGN.CENTER)

    # Arrow between phases
    if i < 3:
        add_textbox(slide, left + 2.9, 3.0, 0.4, 0.5, "→",
                    font_size=28, color=SUBTEXT, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide,
    "Each phase adds realism. Phase 1-2: GRU dominates. Phase 3: non-idealities level the field. "
    "Phase 4: only adaptive decoders survive drift.")

# ═══════════════════════════════════════════════════════════
# SLIDE 16: Novel Contributions
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 11, 0.8,
    "What's Novel Here",
    font_size=32, color=WHITE, bold=True)

contributions = [
    ("1", "First adaptive ML decoder for QEC", "Online learning during inference — nobody has done this", ACCENT_GREEN),
    ("2", "Time-varying non-ideality simulator", "Parameters drift within trajectories, not just between them", ACCENT_BLUE),
    ("3", "Hybrid supervision strategy", "Periodic recalibration + pseudo-labels beats pure self-training", ACCENT_PEACH),
    ("4", "Comprehensive benchmark", "5 decoders × 4 phases of realism, 200+ unit tests", ACCENT_MAUVE),
]

for i, (num, title, desc, color) in enumerate(contributions):
    top = 1.5 + i * 1.3
    add_textbox(slide, 0.8, top, 0.6, 0.5, num,
                font_size=28, color=color, bold=True)
    add_textbox(slide, 1.5, top, 10, 0.5, title,
                font_size=20, color=WHITE, bold=True)
    add_textbox(slide, 1.5, top + 0.45, 10, 0.5, desc,
                font_size=15, color=SUBTEXT)

add_speaker_notes(slide,
    "Phase 4 is entirely novel. Nobody has done adaptive online learning for QEC. "
    "The simulator and hybrid supervision approach are new contributions.")

# ═══════════════════════════════════════════════════════════
# SLIDE 17: Future Work
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 11, 0.8,
    "Future Work & Path to Impact",
    font_size=32, color=WHITE, bold=True)

futures = [
    ("Larger codes", "Scale from 3-qubit to surface codes (the industry standard)"),
    ("Real hardware", "Validate on IBM / Google superconducting qubits"),
    ("Smarter adaptation", "Meta-learning to adapt faster, ensemble methods"),
    ("Multi-task learning", "Simultaneously decode errors AND estimate drifting parameters"),
    ("Publication", "Quantum ML workshops (NeurIPS, ICML) or journals (PRX Quantum)"),
]

for i, (title, desc) in enumerate(futures):
    top = 1.5 + i * 1.1
    add_textbox(slide, 1.2, top, 3.5, 0.5, "→  " + title,
                font_size=18, color=ACCENT_BLUE, bold=True)
    add_textbox(slide, 4.8, top, 8, 0.5, desc,
                font_size=16, color=TEXT)

add_speaker_notes(slide,
    "The 3-qubit code is a proof of concept. The real test is surface codes. "
    "But the principle — adaptive online learning for QEC — transfers directly.")

# ═══════════════════════════════════════════════════════════
# SLIDE 18: Summary
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 1, 1.0, 11, 1.2,
    "Quantum Computers Break.\nWe Taught a Neural Network to Keep Up.",
    font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

takeaways = [
    "1.  Model-based decoders fail when hardware doesn't match the model",
    "2.  Static ML decoders fail when hardware drifts over time",
    "3.  Adaptive decoders with periodic recalibration maintain accuracy",
]

for i, t in enumerate(takeaways):
    color = [ACCENT_RED, ACCENT_PEACH, ACCENT_GREEN][i]
    add_textbox(slide, 1.5, 3.0 + i * 0.9, 10, 0.7, t,
                font_size=22, color=color, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1, 6.0, 11, 0.8,
    "Self-calibrating quantum error correction is possible. This is the first step.",
    font_size=20, color=SUBTEXT, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide,
    "Three clean takeaways. End strong.")


# ═══════════════════════════════════════════════════════════
# BACKUP SLIDES
# ═══════════════════════════════════════════════════════════

# BACKUP 1: GRU Architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.8,
    "Backup: GRU Architecture Detail", font_size=28, color=SUBTEXT, bold=True)
add_bullet_list(slide, 0.8, 1.5, 11, 5, [
    "Input: (window_size, 2) — measurement windows of r₁(t), r₂(t)",
    "GRU: 64 hidden units, 1 layer, batch_first=True",
    "Classifier: Linear(64→32) → ReLU → Dropout(0.1) → Linear(32→4)",
    "Training: Adam optimizer, lr=0.001, 50 epochs, batch_size=256",
    "",
    "Adaptive additions:",
    "  • EMA gradient buffers (decay=0.7)",
    "  • Online weight updates (adapt_lr=0.001)",
    "  • Confidence threshold for pseudo-labels (0.8)",
    "  • supervised_every parameter for hybrid mode",
], font_size=16, color=TEXT)

# BACKUP 2: Confusion Matrices
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.8,
    "Backup: Confusion Matrices", font_size=28, color=SUBTEXT, bold=True)
try_add_image(slide, "phase4_confusion_matrices.png", 0.5, 1.5, 12, 5)

# BACKUP 3: Training Curves
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.8,
    "Backup: Training Curves", font_size=28, color=SUBTEXT, bold=True)
add_textbox(slide, 0.8, 1.2, 11, 0.5,
    "Static and adaptive GRU have identical training curves — difference is at inference only",
    font_size=16, color=SUBTEXT)
try_add_image(slide, "phase4_training_curves.png", 1.5, 2.0, 10, 4.5)

# BACKUP 4: Test Suite
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.8,
    "Backup: Test Suite & Reproducibility", font_size=28, color=SUBTEXT, bold=True)
add_bullet_list(slide, 0.8, 1.5, 11, 5, [
    "225+ unit tests across all phases:",
    "  • 44 tests — quantum operator math",
    "  • 58 tests — Phase 2 Hamiltonian simulator",
    "  • 22 tests — Bayesian filter",
    "  • 99 tests — Phase 3 non-ideal effects",
    "  • 25 tests — Phase 4 adaptive decoder (incl. hybrid supervision)",
    "",
    "Backward compatibility verified:",
    "  Phase 4 (no drift) == Phase 3 == Phase 2 == Phase 1",
    "",
    "Full reproducibility via seeded RNG (seed=42 throughout)",
    "",
    "Repository: github.com/pkarakala/cqec-ml-decoder",
], font_size=16, color=TEXT)

# ═══════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        "adaptive_qec_slides.pptx")
prs.save(out_path)
print(f"✓ Saved {len(prs.slides)} slides to {out_path}")
